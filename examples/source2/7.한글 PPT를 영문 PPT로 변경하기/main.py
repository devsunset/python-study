from pptx import Presentation 
from pptx.util import Pt
import googletrans 
import os
os.chdir(os.path.dirname(os.path.abspath(__file__)))

translator = googletrans.Translator()

str1 = "행복하세요"
result1 = translator.translate(str1, dest='en', src='auto')
print(f"행복하세요 => {result1.text}")

str2 = "I am happy"
result2 = translator.translate(str2, dest='ko', src='en')
print(f"I am happy => {result2.text}")

prs = Presentation("한글.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == 14 or shape.shape_type == 17:
            result = translator.translate(shape.text, dest='en', src='ko') 
            print(shape.text)
            print(result.text)
            shape.text_frame.text = result.text

prs.save("영문변환.pptx")