from pptx import Presentation 
from pptx.util import Pt
import pandas as pd
import copy
import os
os.chdir(os.path.dirname(os.path.abspath(__file__)))

df_from_excel = pd.read_excel("수험번호.xlsx")
prs = Presentation("수험표_샘플.pptx")

이름 = iter(df_from_excel['이름'])
수험번호 = iter(df_from_excel['수험번호'])
응시자수 = len(df_from_excel)

def copy_slide(prs, org_slide):
    copied_slide = prs.slides.add_slide(org_slide.slide_layout)
    for shape in org_slide.shapes:
        org_el = shape.element
        new_el = copy.deepcopy(org_el)
        copied_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    
    for value in org_slide.part.rels:
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.get_or_add(
                value.reltype,
                value._target
                )
    return copied_slide

for i in range(int(응시자수/2)):
    next_slide = copy_slide(prs, prs.slides[0]) 
    for shape in next_slide.shapes:
        if shape.shape_type == 17 and shape.text == "이름":
            shape.text_frame.paragraphs[0].text = next(이름)
            shape.text_frame.paragraphs[0].font.size = Pt(40)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.name = "나눔고딕"
        if shape.shape_type == 17 and shape.text == "수험번호":
            shape.text_frame.paragraphs[0].text = next(수험번호)
            shape.text_frame.paragraphs[0].font.size = Pt(40)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.name = "나눔고딕"

prs.save("수험표_결과.pptx")