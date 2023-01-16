from pptx import Presentation 
from pptx.util import Pt
import os
os.chdir(os.path.dirname(os.path.abspath(__file__)))

prs = Presentation("한글.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == 14 or shape.shape_type == 17:
            print(shape.text)